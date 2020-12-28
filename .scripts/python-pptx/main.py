
# -*- coding:utf-8 -*-
"""
Created in 20XX/XX/XX
@atthor : tajimamasa

 #
 #
 #
 #

"""


from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor


class mainAPP():
    template_path = "./template.pptx"
    save_path = "./output.pptx"    
    presentation = None

    def __init__(self):
        self.presentation = Presentation(self.template_path)
        self.add_title_slide(title="Hello, World!",subtitle= "python-pptx was here!")
        ### 1枚目
        slide = self.add_slide()
        self.slide_title(slide,'no titile')
        self.add_image(slide,"./graph.png")        
        self.add_textbox(slide)

        data = [ ["1.1","1.2","1.3"]
                ,["2.1","2.2","2.3"]
                ,["3.1","3.2","3.3"]]

        table = self.add_paragraph(slide,data)
        self.adjust_table(table)
        
        ### 保存
        self.save()

    ### ファイル保存
    def save(self):
        self.presentation.save(self.save_path)

    def add_title_slide(self,title,subtitle):
        title_slide_layout = self.presentation.slide_layouts[0]
        slide = self.presentation.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    ### 空のスライド作成
    def add_slide(self):
        slide_layout = self.presentation.slide_layouts[1]  # レイアウトや書式を元ファイルから参照する
        slide = self.presentation.slides.add_slide(slide_layout)
        return slide

    ### スライドタイトル記入
    def slide_title(self,slide,slide_title):
        slide.shapes.title.text = slide_title

    ### imageの挿入
    def add_image(self,slide,image_path,pos=[1,5],picheight=[7.9]):
        pic_left = self.Centis(pos[0])
        pic_top = self.Centis(pos[1])
        pic_height = self.Centis(picheight[0])
        slide.shapes.add_picture(image_path, pic_left, pic_top, height=pic_height)


    ### テキストボックスを追加
    def add_textbox(self,slide,text="Test Text",fontsize=20,pos=[1,13.4],size=[25.4-2,5]):
        text_left = self.Centis(pos[0])
        text_top = self.Centis(pos[1])

        text_width = self.Centis(size[0])
        text_height = self.Centis(size[1])

        color = RGBColor(255, 255, 255)

        text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)

        text_box.text = text
        text_box.text_frame.add_paragraph().font.size = Pt(fontsize)
        text_box.fill.solid()
        text_box.fill.fore_color.rgb = color

        
    ### センチメートルをPowerPoint上の距離に変換する関数
    def Centis(self,length):
        centi = Inches(length/2.54)
        return centi

    ### 表を挿入
    def add_paragraph(self,slide,data,fontsize=20,pos=[9.4,5]):

        table_left = self.Centis(pos[0])
        table_top = self.Centis(pos[1])

        # tableの幅と高さ（仮）
        table_width = self.Centis(15)
        table_height = self.Centis(10)

        # tableの行数と列数(tableのサイズ)
        rows = len(data)
        cols = len(data[0])

        table = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height).table

        # 表の各セルの中身を記入
        for i in range(rows):
            for j in range(cols):
                cell = table.cell(i, j)
                cell.text = data[i][j]
                cell.text_frame.paragraphs[0].font.size = Pt(fontsize)
        return table

    ### 表を再調整
    def adjust_table(self,table):
        # tableの高さを再調整
        table.rows[0].height = self.Centis(1.5)
        table.rows[1].height = self.Centis(4.9)
        table.rows[2].height = self.Centis(1.5)

        # tableの幅を再調整
        table.columns[0].width = self.Centis((15) / 3)
        table.columns[1].width = self.Centis((15) / 3)
        table.columns[2].width = self.Centis((15) / 3)


if __name__=='__main__':
    mainAPP()
